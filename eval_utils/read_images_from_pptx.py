import re
import os
from io import BytesIO
from pptx import Presentation
from PIL import Image

images_folder = 'images/'
os.makedirs(images_folder, exist_ok=True)

prs = Presentation('correct_object_info.pptx')
for s_idx, slide in enumerate(prs.slides):
    s_idx += 1
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text = shape.text.lower().strip()
            print(text)
            text = re.sub(r'&', '_and_', text)
            text = re.sub('\'', '', text)
            text = re.sub('\s', '_', text)
            print(text)
        elif hasattr(shape, 'image'):
            image = shape.image
            image_bytes = BytesIO(image.blob)
            img = Image.open(image_bytes)

            image_type = image.content_type.split('/')[-1]
            img_filename = os.path.join(images_folder, f"{s_idx:03d}_{text}.{image_type}")
            print(img_filename)
            img.save(img_filename)

